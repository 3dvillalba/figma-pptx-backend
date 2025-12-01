#!/usr/bin/env python3
"""
generate_pptx.py - Genera PPTX VERTICAL sincronizado con Figma stage
"""
import json
import sys
import base64
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_vertical_pptx(slides_data):
    """
    Crea PPTX sincronizado con dimensiones Figma
    
    Cada slide recibe ancho/alto en píxeles Figma (72 DPI)
    Se convierte automáticamente a pulgadas para PowerPoint
    """
    
    prs = Presentation()
    
    print(f"Creando {len(slides_data)} slides sincronizados con Figma...")
    
    for idx, slide_data in enumerate(slides_data):
        # Obtener dimensiones del frame Figma (en píxeles 72 DPI)
        figma_width_px = slide_data.get('width', 960)
        figma_height_px = slide_data.get('height', 1280)
        
        # Convertir píxeles Figma (72 DPI) a pulgadas
        # Fórmula: Pulgadas = Píxeles / 72
        width_inches = figma_width_px / 72.0
        height_inches = figma_height_px / 72.0
        
        print(f"Slide {idx + 1}: {figma_width_px}px × {figma_height_px}px = {width_inches:.2f}\" × {height_inches:.2f}\"")
        
        # SINCRONIZAR: Aplicar dimensiones exactas a PowerPoint
        prs.slide_width = Inches(width_inches)
        prs.slide_height = Inches(height_inches)
        
        # Crear slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Fondo blanco
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Agregar imagen a pantalla completa (sincronizada)
        elements = slide_data.get('elements', [])
        for element in elements:
            if element.get('type') == 'image':
                add_image_fullscreen(slide, element, width_inches, height_inches)
    
    return prs

def add_image_fullscreen(slide, element, width_inches, height_inches):
    """Agrega imagen escalada a dimensiones completas del slide"""
    
    image_base64 = element.get('imageBase64')
    if not image_base64:
        return
    
    try:
        # Decodificar base64
        image_bytes = base64.b64decode(image_base64)
        image_stream = io.BytesIO(image_bytes)
        
        # SINCRONIZAR: Imagen ocupa 100% del slide
        # left=0, top=0, width/height = dimensiones slide exactas
        slide.shapes.add_picture(
            image_stream,
            left=Inches(0),
            top=Inches(0),
            width=Inches(width_inches),
            height=Inches(height_inches)
        )
        
        print(f"  ✓ Imagen agregada {width_inches:.2f}\" × {height_inches:.2f}\"")
        
    except Exception as e:
        print(f"  ✗ Error agregando imagen: {str(e)}")

if __name__ == '__main__':
    if len(sys.argv) != 3:
        print("Uso: python3 generate_pptx.py <input_json> <output_pptx>")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    
    try:
        with open(input_file, 'r') as f:
            data = json.load(f)
        
        slides = data.get('slides', [])
        prs = create_vertical_pptx(slides)
        prs.save(output_file)
        
        print(f"\n✓ PPTX sincronizado con Figma generado: {output_file}")
        
    except Exception as e:
        print(f"ERROR: {str(e)}", file=sys.stderr)
        sys.exit(1)
